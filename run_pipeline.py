import argparse
import os
from src.summarizer import extract_sections_from_pdf
from src.video_maker import make_video
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor


def create_slides(sections, outdir="outputs"):
    """
    Generate PowerPoint slides from extracted text sections.
    """
    os.makedirs(outdir, exist_ok=True)
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]

    for i, section in enumerate(sections, start=1):
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]

        title_text = section.get("title", f"Section {i}")
        content_text = section.get("content", "No content available.")

        title_shape.text = title_text
        content_shape.text = content_text

        # Formatting
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(40, 40, 40)

        for p in content_shape.text_frame.paragraphs:
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(60, 60, 60)

    pptx_path = os.path.join(outdir, "slides.pptx")
    prs.save(pptx_path)
    print(f"📊 Slides generated successfully → {pptx_path}")
    return pptx_path


def main():
    parser = argparse.ArgumentParser(description="PDF → Slides → Video automation")
    parser.add_argument("--input", required=True, help="Input PDF file")
    parser.add_argument("--outdir", default="outputs", help="Output folder")
    args = parser.parse_args()

    pdf_path = args.input
    os.makedirs(args.outdir, exist_ok=True)

    print(f"📂 Using input file: {os.path.basename(pdf_path)}")

    # Step 1: Extract text
    sections = extract_sections_from_pdf(pdf_path)
    if not sections:
        print("⚠️ No text found in PDF.")
        return

    # Step 2: Create slides
    create_slides(sections, args.outdir)

    # Step 3: Generate video
    make_video(sections, outdir=args.outdir)

    print("✅ Done! Check your outputs folder for slides.pptx and video.mp4 🎬")


if __name__ == "__main__":
    main()
