from pptx import Presentation
from pptx.util import Inches

def build_pptx(slides, outdir='outputs'):
    prs = Presentation()
    for s in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = s['title']
        left, top, w, h = Inches(1), Inches(1.5), Inches(8), Inches(4.5)
        slide.shapes.add_picture(s['image_path'], left, top, width=w, height=h)
        body_shape = slide.placeholders[1].text_frame
        for b in s['bullets']:
            p = body_shape.add_paragraph()
            p.text = b
    prs.save(f"{outdir}/slides.pptx")
